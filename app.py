import io
import os
import platform
import re
import subprocess
import tempfile
from pathlib import Path
from urllib.parse import urlparse

import requests
import streamlit as st
from bs4 import BeautifulSoup
from google import genai
from google.genai import types

st.set_page_config(page_title="Study Agent", page_icon="📚", layout="wide")

MODEL = "gemini-2.5-flash"
PRICING_SUBJECT = "💰 Pricing (MBA)"

# ── System Prompts ────────────────────────────────────────────────────────────

SUMMARY_SYSTEM_PROMPT = """You are an expert study assistant. Process academic materials and present knowledge in a clear, structured way.

Structure your output as:

## Overview
2-3 sentence summary of what this material covers.

## Key Concepts
The most important ideas, frameworks, and theories — explained clearly.

## Important Details
Supporting facts, examples, data, or evidence worth remembering.

## Connections
How these ideas connect to broader principles or other topics.

## What to Remember
3-5 bullet points of the most high-leverage takeaways.

Adapt your depth and tone to the subject matter."""

PODCAST_SYSTEM_PROMPT = """You are an engaging podcast host who makes complex academic material come alive through storytelling.

- Open with a vivid hook that draws the listener in
- Tell the human story BEHIND the concepts — who developed these ideas, what problem they solved
- Explain frameworks and theories through narrative, not bullet points
- Use vivid analogies that will stick in memory
- Connect the material to a broader principle at the end
- Sound like a brilliant friend explaining over coffee, not a professor lecturing
- End with 3 crisp memorable takeaways

Format as a natural podcast script (5–10 minutes of speaking time). Use paragraph breaks for natural pauses. \
Mark key terms in [BRACKETS] when first introduced."""

PRICING_PODCAST_SYSTEM_PROMPT = """You are "The Pricing Professor" — an engaging podcast host who makes complex \
pricing strategies and business cases come alive through storytelling.

- Open with a vivid hook: "Picture this: It's 1998, and a small airline is about to change how the world buys tickets..."
- Tell the human story BEHIND the pricing strategy — who made the decision, what pressure they were under
- Explain frameworks (price discrimination, value-based pricing, etc.) through narrative, not bullet points
- Use vivid analogies that MBA students will remember on exam day
- Connect the case to a broader principle at the end
- Sound like a brilliant friend explaining over coffee, not a professor lecturing
- End with 3 crisp takeaways framed as memorable insights

Format as a natural podcast script (5–10 minutes of speaking time). Use paragraph breaks for natural pauses. \
Mark key terms in [BRACKETS] when first introduced."""

QA_SYSTEM_PROMPT = """You are a sharp professor helping a student study. You have access to their course materials.

- Answer questions directly, then add depth
- Connect concepts across multiple materials when relevant
- Challenge the student with follow-up questions when appropriate
- Use real-world examples to reinforce concepts
- Flag what's likely to appear on an exam vs. what's deeper context"""

PRICING_QA_SYSTEM_PROMPT = """You are a sharp MBA pricing professor helping a student prepare for exams. \
You have access to their case studies and slides.

- Answer questions directly, then add depth
- Connect concepts across multiple cases when relevant
- Challenge the student with follow-up questions when appropriate
- Use real-world examples to reinforce concepts
- Flag what's likely to appear on an exam vs. what's deeper context"""


# ── Session State ─────────────────────────────────────────────────────────────

def init_session():
    defaults = {
        "api_key": os.environ.get("GEMINI_API_KEY", ""),
        "subjects": {PRICING_SUBJECT: {"files": [], "urls": []}},
        "active_subject": PRICING_SUBJECT,
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

init_session()


# ── Helpers ───────────────────────────────────────────────────────────────────

def get_client():
    api_key = st.session_state.get("api_key") or os.environ.get("GEMINI_API_KEY", "")
    if not api_key:
        st.error("Enter your Gemini API key in the sidebar to get started.")
        st.stop()
    return genai.Client(api_key=api_key)


def fetch_url(url: str) -> str:
    """Fetch and extract readable text from a webpage."""
    headers = {"User-Agent": "Mozilla/5.0 (compatible; StudyAgent/1.0)"}
    r = requests.get(url, headers=headers, timeout=15)
    r.raise_for_status()
    soup = BeautifulSoup(r.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    lines = [l for l in text.splitlines() if l.strip()]
    return "\n".join(lines)


def extract_text(file_bytes: bytes, filename: str):
    """Extract plain text from docx/pptx. Returns None for native types."""
    suffix = Path(filename).suffix.lower()
    if suffix in (".docx", ".doc"):
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        lines = [
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        return "\n".join(lines)
    return None


def upload_to_gemini(client, file_bytes: bytes, filename: str):
    text = extract_text(file_bytes, filename)
    suffix = ".txt" if text else Path(filename).suffix.lower()
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(text.encode() if text else file_bytes)
        tmp_path = tmp.name
    uploaded = client.files.upload(file=tmp_path, config={"display_name": filename})
    Path(tmp_path).unlink(missing_ok=True)
    return uploaded


def upload_text_to_gemini(client, text: str, display_name: str):
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False, mode="w", encoding="utf-8") as tmp:
        tmp.write(text)
        tmp_path = tmp.name
    uploaded = client.files.upload(file=tmp_path, config={"display_name": display_name})
    Path(tmp_path).unlink(missing_ok=True)
    return uploaded


def speak(text: str):
    if platform.system() != "Darwin":
        return
    clean = re.sub(r"\*+|_+|#{1,6}\s*|`+|\[|\]", "", text)
    subprocess.Popen(["say", "-v", "Samantha", "-r", "175", clean])


def cleanup_gemini_files(client, files: list):
    for f in files:
        try:
            client.files.delete(name=f.name)
        except Exception:
            pass


def gather_parts(client, subject_data: dict, status=None):
    """Upload all files and URLs for a subject. Returns (parts, gemini_files)."""
    parts = []
    gemini_files = []
    for item in subject_data.get("files", []):
        if status:
            status.write(f"⬆️ {item['name']}")
        gf = upload_to_gemini(client, item["bytes"], item["name"])
        gemini_files.append(gf)
        parts.append(gf)
    for item in subject_data.get("urls", []):
        if status:
            status.write(f"🌐 Fetching {item['url']}...")
        try:
            text = fetch_url(item["url"])
            gf = upload_text_to_gemini(client, text, item["url"])
            gemini_files.append(gf)
            parts.append(gf)
        except Exception as e:
            if status:
                status.write(f"⚠️ Could not fetch {item['url']}: {e}")
    return parts, gemini_files


# ── Sidebar ───────────────────────────────────────────────────────────────────

with st.sidebar:
    st.title("📚 Study Agent")
    st.caption("AI-powered study companion")
    st.divider()

    key_input = st.text_input(
        "Gemini API Key",
        value=st.session_state.api_key,
        type="password",
        help="Get your key at aistudio.google.com/app/apikey",
    )
    if key_input:
        st.session_state.api_key = key_input

    st.divider()
    speak_aloud = st.toggle("🔊 Read aloud", help="macOS built-in TTS (Samantha voice)")
    st.divider()

    st.markdown("**Subjects**")
    new_sub = st.text_input("New subject", placeholder="Economics, History, CS...")
    if st.button("➕ Add subject", disabled=not new_sub.strip()):
        name = new_sub.strip()
        if name not in st.session_state.subjects:
            st.session_state.subjects[name] = {"files": [], "urls": []}
        st.session_state.active_subject = name
        st.rerun()

    st.divider()
    for sub in list(st.session_state.subjects.keys()):
        n_files = len(st.session_state.subjects[sub]["files"])
        n_urls = len(st.session_state.subjects[sub]["urls"])
        badge = f" · {n_files}📄 {n_urls}🔗" if (n_files or n_urls) else ""
        is_active = sub == st.session_state.active_subject
        if st.button(
            f"{sub}{badge}",
            key=f"sub_btn_{sub}",
            type="primary" if is_active else "secondary",
            use_container_width=True,
        ):
            st.session_state.active_subject = sub
            st.rerun()


# ── Main Area ─────────────────────────────────────────────────────────────────

subject = st.session_state.active_subject
subject_data = st.session_state.subjects[subject]
is_pricing = subject == PRICING_SUBJECT
has_materials = bool(subject_data["files"] or subject_data["urls"])

st.title(subject)

# ── Add Materials ─────────────────────────────────────────────────────────────

with st.expander("📁 Add Materials", expanded=not has_materials):
    col_file, col_url = st.columns(2)

    with col_file:
        st.markdown("**Upload files**")
        new_files = st.file_uploader(
            "Upload",
            type=["pdf", "docx", "doc", "pptx", "txt", "md"],
            accept_multiple_files=True,
            key=f"uploader_{subject}",
            label_visibility="collapsed",
        )
        if new_files and st.button("➕ Add files", key=f"add_files_{subject}"):
            existing = {f["name"] for f in subject_data["files"]}
            added = 0
            for f in new_files:
                if f.name not in existing:
                    subject_data["files"].append({"name": f.name, "bytes": f.read()})
                    added += 1
            if added:
                st.success(f"Added {added} file(s)")
                st.rerun()

    with col_url:
        st.markdown("**Add a URL**")
        url_input = st.text_input(
            "URL",
            placeholder="https://...",
            key=f"url_input_{subject}",
            label_visibility="collapsed",
        )
        if st.button("➕ Add URL", key=f"add_url_{subject}", disabled=not url_input.strip()):
            existing = {u["url"] for u in subject_data["urls"]}
            if url_input not in existing:
                subject_data["urls"].append({
                    "url": url_input,
                    "title": urlparse(url_input).netloc or url_input,
                })
                st.success("URL added")
                st.rerun()

# Current materials list
if has_materials:
    with st.expander(
        f"📚 Materials — {len(subject_data['files'])} file(s), {len(subject_data['urls'])} URL(s)",
        expanded=False,
    ):
        for i, f in enumerate(list(subject_data["files"])):
            c1, c2 = st.columns([8, 1])
            c1.caption(f"📄 {f['name']}")
            if c2.button("✕", key=f"rm_file_{subject}_{i}"):
                subject_data["files"].pop(i)
                st.rerun()
        for i, u in enumerate(list(subject_data["urls"])):
            c1, c2 = st.columns([8, 1])
            c1.caption(f"🔗 {u['url']}")
            if c2.button("✕", key=f"rm_url_{subject}_{i}"):
                subject_data["urls"].pop(i)
                st.rerun()

if not has_materials:
    st.info("Add files or URLs above, then choose a mode below.")
    st.stop()

# ── Tabs ──────────────────────────────────────────────────────────────────────

summary_tab, podcast_tab, qa_tab = st.tabs(["📋 Summary", "🎙️ Podcast", "💬 Q&A"])

# ── Summary Tab ───────────────────────────────────────────────────────────────

with summary_tab:
    st.caption("Structured breakdown: key concepts, frameworks, and what to remember.")
    focus = st.text_input(
        "Focus topic (optional)",
        placeholder="e.g. supply and demand, price discrimination",
        key="sum_focus",
    )
    if st.button("📋 Generate Summary", type="primary", key="gen_summary"):
        client = get_client()
        with st.status("Processing materials...", expanded=True) as status:
            parts, gemini_files = gather_parts(client, subject_data, status)
            if not parts:
                st.error("No materials could be loaded.")
                st.stop()
            status.update(label="Generating summary...", state="running")
        try:
            focus_note = f"\n\nFocus especially on: {focus}" if focus else ""
            prompt = f"Create a structured study summary of these materials.{focus_note}"
            full_text = ""
            out = st.empty()
            for chunk in client.models.generate_content_stream(
                model=MODEL,
                contents=parts + [prompt],
                config=types.GenerateContentConfig(system_instruction=SUMMARY_SYSTEM_PROMPT),
            ):
                full_text += chunk.text or ""
                out.markdown(full_text)
            st.download_button(
                "💾 Download summary",
                data=full_text,
                file_name=f"summary_{subject}.txt",
                mime="text/plain",
            )
            if speak_aloud:
                speak(full_text)
        finally:
            cleanup_gemini_files(client, gemini_files)

# ── Podcast Tab ───────────────────────────────────────────────────────────────

with podcast_tab:
    st.caption("Story-style narration — understand your materials through storytelling.")
    focus = st.text_input(
        "Focus topic (optional)",
        placeholder="e.g. Keynesian economics, price discrimination",
        key="pod_focus",
    )
    if st.button("🎙️ Generate Podcast", type="primary", key="gen_podcast"):
        client = get_client()
        with st.status("Processing materials...", expanded=True) as status:
            parts, gemini_files = gather_parts(client, subject_data, status)
            if not parts:
                st.error("No materials could be loaded.")
                st.stop()
            status.update(label="Generating narration...", state="running")
        try:
            focus_note = f"\n\nPay special attention to: {focus}" if focus else ""
            prompt = f"Narrate this as a podcast episode for a student studying this subject.{focus_note}"
            system = PRICING_PODCAST_SYSTEM_PROMPT if is_pricing else PODCAST_SYSTEM_PROMPT
            full_text = ""
            out = st.empty()
            for chunk in client.models.generate_content_stream(
                model=MODEL,
                contents=parts + [prompt],
                config=types.GenerateContentConfig(system_instruction=system),
            ):
                full_text += chunk.text or ""
                out.markdown(full_text)
            st.download_button(
                "💾 Download transcript",
                data=full_text,
                file_name=f"podcast_{subject}.txt",
                mime="text/plain",
            )
            if speak_aloud:
                speak(full_text)
        finally:
            cleanup_gemini_files(client, gemini_files)

# ── Q&A Tab ───────────────────────────────────────────────────────────────────

with qa_tab:
    st.caption("Chat with an AI professor about your materials.")

    qa_key = f"qa_{subject}"
    for k, d in [
        (f"{qa_key}_messages", []),
        (f"{qa_key}_gemini_files", []),
        (f"{qa_key}_chat", None),
        (f"{qa_key}_ready", False),
    ]:
        if k not in st.session_state:
            st.session_state[k] = d

    if not st.session_state[f"{qa_key}_ready"]:
        if st.button("🚀 Start Study Session", type="primary"):
            client = get_client()
            with st.status("Loading materials...", expanded=True) as status:
                parts, gemini_files = gather_parts(client, subject_data, status)
                if not parts:
                    st.error("No materials could be loaded.")
                    st.stop()
                status.update(label="Ready!", state="complete")

            st.session_state[f"{qa_key}_gemini_files"] = gemini_files
            system = PRICING_QA_SYSTEM_PROMPT if is_pricing else QA_SYSTEM_PROMPT
            st.session_state[f"{qa_key}_chat"] = client.chats.create(
                model=MODEL,
                config=types.GenerateContentConfig(system_instruction=system),
            )
            initial = parts + [
                "I've uploaded my study materials. "
                "Acknowledge what you have and offer 3 good study questions to start."
            ]
            intro = "".join(
                chunk.text or ""
                for chunk in st.session_state[f"{qa_key}_chat"].send_message_stream(initial)
            )
            st.session_state[f"{qa_key}_messages"] = [{"role": "assistant", "content": intro}]
            st.session_state[f"{qa_key}_ready"] = True
            if speak_aloud:
                speak(intro)
            st.rerun()

    else:
        if st.button("🔄 New Session"):
            client = get_client()
            cleanup_gemini_files(client, st.session_state[f"{qa_key}_gemini_files"])
            st.session_state[f"{qa_key}_ready"] = False
            st.session_state[f"{qa_key}_messages"] = []
            st.session_state[f"{qa_key}_gemini_files"] = []
            st.session_state[f"{qa_key}_chat"] = None
            st.rerun()

        for msg in st.session_state[f"{qa_key}_messages"]:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

        if user_input := st.chat_input("Ask anything about your materials..."):
            st.session_state[f"{qa_key}_messages"].append({"role": "user", "content": user_input})
            with st.chat_message("user"):
                st.markdown(user_input)
            with st.chat_message("assistant"):
                response_text = ""
                placeholder = st.empty()
                for chunk in st.session_state[f"{qa_key}_chat"].send_message_stream(user_input):
                    response_text += chunk.text or ""
                    placeholder.markdown(response_text)
            st.session_state[f"{qa_key}_messages"].append({"role": "assistant", "content": response_text})
            if speak_aloud:
                speak(response_text)
