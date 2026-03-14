#!/usr/bin/env python3
"""
Pricing Class Podcast Agent
============================
Converts your MBA pricing case studies and slides into:
  - Engaging podcast-style narrations (storytelling mode)
  - Interactive Q&A study sessions

Usage:
  python3 podcast_agent.py podcast amazon
  python3 podcast_agent.py podcast amazon uber --focus "price discrimination"
  python3 podcast_agent.py qa amazon uber
"""

import os
import re
import subprocess
import sys
from google import genai
from google.genai import types
from pathlib import Path

client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])

MODEL = "gemini-2.5-flash"

PODCAST_SYSTEM_PROMPT = """You are "The Pricing Professor" — an engaging podcast host who makes complex \
pricing strategies and business cases come alive through storytelling.

Your style:
- Open with a vivid hook: "Picture this: It's 1998, and a small airline is about to change how \
the world buys tickets..."
- Tell the human story BEHIND the pricing strategy — who made the decision, what pressure they were under
- Explain frameworks (price discrimination, value-based pricing, etc.) through narrative, not bullet points
- Use vivid analogies that MBA students will remember on exam day
- Connect the case to a broader principle at the end
- Sound like a brilliant friend explaining over coffee, not a professor lecturing
- End with 3 crisp takeaways framed as memorable insights, not dry summaries

Format as a natural podcast script (5–10 minutes of speaking time). Use paragraph breaks for \
natural pauses. Mark key terms in [BRACKETS] when first introduced."""

QA_SYSTEM_PROMPT = """You are a sharp MBA pricing professor helping a student prepare for exams and \
deepen their understanding. You have access to their case studies and slides.

Your style:
- Answer questions directly, then add depth
- Connect concepts across multiple cases when relevant
- Challenge the student with follow-up questions when appropriate
- Use real-world examples to reinforce concepts
- Flag what's likely to appear on an exam vs. what's deeper context"""

SEARCH_DIRS = [
    Path.cwd(),
    Path.home() / "pricing_study",
    Path.home() / "Downloads",
    Path.home() / "Desktop",
    Path.home() / "Documents",
]


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".txt", ".md", ".pptx"}

MIME_MAP = {
    ".pdf":  "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc":  "application/msword",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt":  "text/plain",
    ".md":   "text/plain",
}


def find_files(query: str) -> list:
    """Find supported files matching a company/keyword across common folders."""
    query_lower = query.lower()
    matches = []
    seen = set()
    for folder in SEARCH_DIRS:
        if not folder.exists():
            continue
        for f in folder.iterdir():
            if f.suffix.lower() in SUPPORTED_EXTENSIONS and query_lower in f.name.lower() and f not in seen:
                matches.append(f)
                seen.add(f)
    return matches


def resolve_files(args: list) -> list:
    """Turn company names or file paths into resolved file paths."""
    resolved = []
    for arg in args:
        path = Path(arg)
        if path.exists() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            resolved.append(str(path))
            continue
        matches = find_files(arg)
        if not matches:
            searched = [str(d) for d in SEARCH_DIRS if d.exists()]
            print(f"  No PDFs found matching '{arg}' in: {searched}")
            sys.exit(1)
        if len(matches) == 1:
            print(f"  Found: {matches[0].name}")
            resolved.append(str(matches[0]))
        else:
            print(f"\n  Multiple matches for '{arg}':")
            for i, m in enumerate(matches, 1):
                print(f"    {i}. {m.name}  ({m.parent})")
            choice = input("  Pick a number: ").strip()
            if not choice.isdigit() or not (1 <= int(choice) <= len(matches)):
                print("  Invalid choice.")
                sys.exit(1)
            resolved.append(str(matches[int(choice) - 1]))
    return resolved


def extract_text(file_path: str) -> str:
    """Extract plain text from docx or pptx files."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix in (".docx", ".doc"):
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)
    return None  # Use native upload for PDF/txt/md


def upload_file(file_path: str):
    """Upload a file to Gemini. Converts docx/pptx to text first if needed."""
    path = Path(file_path)
    print(f"  Uploading {path.name}...")

    text = extract_text(file_path)
    if text:
        # Upload as plain text
        import tempfile
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as tmp:
            tmp.write(text)
            tmp_path = tmp.name
        uploaded = client.files.upload(file=tmp_path, config={"display_name": path.name})
        Path(tmp_path).unlink()
    else:
        uploaded = client.files.upload(file=file_path)

    print(f"  Done: {uploaded.name}")
    return uploaded


def speak(text: str):
    """Read text aloud using macOS built-in TTS. Strips markdown formatting."""
    clean = re.sub(r"\*+|_+|#{1,6}\s*|`+|\[|\]", "", text)
    subprocess.run(["say", "-v", "Samantha", "-r", "175", clean])


def narrate_as_podcast(gemini_file, focus: str = None, aloud: bool = False) -> str:
    """Stream a podcast-style narration of a case study."""
    focus_note = f"\n\nPay special attention to: {focus}" if focus else ""
    print(f"\n  {gemini_file.display_name or gemini_file.name}")
    print("  " + "-" * 56)

    prompt = f"Narrate this as a podcast episode for an MBA pricing student.{focus_note}"
    full_text = ""

    for chunk in client.models.generate_content_stream(
        model=MODEL,
        contents=[gemini_file, prompt],
        config=types.GenerateContentConfig(system_instruction=PODCAST_SYSTEM_PROMPT),
    ):
        text = chunk.text or ""
        print(text, end="", flush=True)
        full_text += text

    print("\n")
    if aloud:
        print("  Reading aloud... (press Ctrl+C to stop)")
        speak(full_text)
    return full_text


def save_narration(display_name: str, narration: str) -> str:
    """Save narration to a text file. Returns the output path."""
    stem = Path(display_name).stem
    output_path = f"podcast_{stem}.txt"
    with open(output_path, "w") as f:
        f.write(f"# Podcast Narration: {display_name}\n\n")
        f.write(narration)
    return output_path


def podcast_mode(file_paths: list, focus: str = None, aloud: bool = False):
    """Upload files and narrate each one as a podcast episode."""
    print("\nPricing Class Podcast Agent")
    print("=" * 60)
    print("\nUploading your materials...")

    uploaded = [upload_file(fp) for fp in file_paths]

    print("\nGenerating podcast narrations...")
    try:
        for gemini_file in uploaded:
            narration = narrate_as_podcast(gemini_file, focus, aloud=aloud)
            name = gemini_file.display_name or gemini_file.name
            output_path = save_narration(name, narration)
            print(f"  Saved to: {output_path}")
    finally:
        print("\nCleaning up uploaded files...")
        for gemini_file in uploaded:
            client.files.delete(name=gemini_file.name)

    print("\nDone! Your podcast scripts are ready.")


def qa_mode(file_paths: list, aloud: bool = False):
    """Interactive Q&A session about uploaded materials."""
    print("\nPricing Class Q&A Agent")
    print("=" * 60)
    print("\nUploading your materials...")

    uploaded = [upload_file(fp) for fp in file_paths]

    print("\nReady! Ask me anything about your pricing materials.")
    print("Type 'quit' to exit\n")

    chat = client.chats.create(
        model=MODEL,
        config=types.GenerateContentConfig(system_instruction=QA_SYSTEM_PROMPT),
    )

    initial_parts = uploaded + [
        "I've uploaded my MBA pricing case studies and slides. "
        "Please acknowledge what you have and offer 3 good study questions to start."
    ]

    try:
        print("Professor: ", end="", flush=True)
        intro = ""
        for chunk in chat.send_message_stream(initial_parts):
            text = chunk.text or ""
            print(text, end="", flush=True)
            intro += text
        print("\n")
        if aloud:
            speak(intro)

        while True:
            try:
                user_input = input("You: ").strip()
            except (EOFError, KeyboardInterrupt):
                break

            if not user_input:
                continue
            if user_input.lower() in ("quit", "exit", "q"):
                break

            print("Professor: ", end="", flush=True)
            response_text = ""
            for chunk in chat.send_message_stream(user_input):
                text = chunk.text or ""
                print(text, end="", flush=True)
                response_text += text
            print("\n")
            if aloud:
                speak(response_text)

    finally:
        print("\nCleaning up uploaded files...")
        for gemini_file in uploaded:
            client.files.delete(name=gemini_file.name)
        print("Goodbye! Good luck with your pricing class.")


def print_usage():
    print("""
Pricing Class Podcast Agent
============================

Podcast mode (narrate as a story):
  python3 podcast_agent.py podcast amazon
  python3 podcast_agent.py podcast amazon uber --focus "price discrimination"

Q&A mode (interactive study session):
  python3 podcast_agent.py qa amazon uber

Use company names, keywords, or full file paths.
PDFs are searched in: current folder, ~/pricing_study, ~/Downloads, ~/Desktop, ~/Documents

Set your API key first:
  export GEMINI_API_KEY=your-key-here
""")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print_usage()
        sys.exit(0)

    mode = sys.argv[1].lower()
    args = sys.argv[2:]

    aloud = "--speak" in args
    if aloud:
        args = [a for a in args if a != "--speak"]

    focus = None
    if "--focus" in args:
        idx = args.index("--focus")
        if idx + 1 < len(args):
            focus = args[idx + 1]
            args = [a for i, a in enumerate(args) if i not in (idx, idx + 1)]
        else:
            print("  --focus requires a value, e.g. --focus 'price discrimination'")
            sys.exit(1)

    if not args:
        print("  Please provide at least one file name or company name.")
        sys.exit(1)

    file_paths = resolve_files(args)

    if mode == "podcast":
        podcast_mode(file_paths, focus, aloud=aloud)
    elif mode == "qa":
        qa_mode(file_paths, aloud=aloud)
    else:
        print(f"  Unknown mode '{mode}'. Use 'podcast' or 'qa'.")
        print_usage()
        sys.exit(1)
